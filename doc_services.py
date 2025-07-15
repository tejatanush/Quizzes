import nltk
import fitz
import os
import base64
project_root = os.path.dirname(os.path.abspath(__file__))  
nltk_data_path = os.path.join(project_root, "nltk_data")
nltk.data.path.append(nltk_data_path)
import tiktoken
from nltk.tokenize import sent_tokenize
from langchain_core.prompts import PromptTemplate
from langchain_groq import ChatGroq
from langchain_core.output_parsers import StrOutputParser
from io import BytesIO
from PIL import Image
import google.generativeai as genai
import re
import json
from config import GROQ_API_KEY, google_Api_key
import time
import logging
from typing import List, Dict, Any, Tuple
from pathlib import Path
import asyncio
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class QuizGenerationService:
    def __init__(self):
        self.supported_formats = {'.pdf', '.docx', '.pptx'}
        self.max_retries = 3
        self.wait_time = 4
        
    def validate_file(self, filepath: str) -> bool:
        """Validate if file exists and has supported format"""
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"File not found: {filepath}")
        
        file_extension = Path(filepath).suffix.lower()
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: {self.supported_formats}")
        
        return True

    def get_and_extract_file(self, filepath: str) -> Tuple[str, List[str]]:
        """Extract text and images from PDF, DOCX, or PPTX files"""
        self.validate_file(filepath)
        
        file_extension = Path(filepath).suffix.lower()
        
        try:
            if file_extension == '.pdf':
                return self._extract_from_pdf(filepath)
            elif file_extension == '.docx':
                return self._extract_from_docx(filepath)
            elif file_extension == '.pptx':
                return self._extract_from_pptx(filepath)
        except Exception as e:
            logger.error(f"Error extracting from {filepath}: {str(e)}")
            raise Exception(f"Failed to extract content from file: {str(e)}")

    def _extract_from_pdf(self, filepath: str) -> Tuple[str, List[str]]:
        """Extract text and images from PDF"""
        doc = fitz.open(filepath)
        text_content = "\n".join([page.get_text() for page in doc])
        image_content = []
        
        for i, page in enumerate(doc):
            for img_index, img in enumerate(page.get_images(full=True)):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_b64 = base64.b64encode(image_bytes).decode("utf-8")
                    image_content.append(image_b64)
                except Exception as e:
                    logger.warning(f"Failed to extract image {img_index} from page {i}: {str(e)}")
                    continue
        
        doc.close()
        return text_content, image_content

    def _extract_from_docx(self, filepath: str) -> Tuple[str, List[str]]:
        """Extract text and images from DOCX"""
        doc = Document(filepath)
        text_content = ""
        image_content = []
        for paragraph in doc.paragraphs:
            text_content += paragraph.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_content += cell.text + "\t"
                text_content += "\n"
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    image_b64 = base64.b64encode(image_data).decode("utf-8")
                    image_content.append(image_b64)
                except Exception as e:
                    logger.warning(f"Failed to extract image from DOCX: {str(e)}")
                    continue
        
        return text_content, image_content

    def _extract_from_pptx(self, filepath: str) -> Tuple[str, List[str]]:
        """Extract text and images from PPTX"""
        prs = Presentation(filepath)
        text_content = ""
        image_content = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"
                
                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_data = shape.image.blob
                        image_b64 = base64.b64encode(image_data).decode("utf-8")
                        image_content.append(image_b64)
                    except Exception as e:
                        logger.warning(f"Failed to extract image from PPTX slide: {str(e)}")
                        continue
        
        return text_content, image_content

    def split_text_into_chunks(self, text: str, model_name: str = "gpt-3.5-turbo", max_tokens: int = 3000) -> List[str]:
        """Split text into chunks based on token limits"""
        if not text.strip():
            return []
        
        try:
            tokenizer = tiktoken.encoding_for_model(model_name)
            sentences = sent_tokenize(text)
            chunks = []
            current_chunk = ""
            current_tokens = 0

            for sentence in sentences:
                sentence = sentence.strip()
                if not sentence:
                    continue

                sentence_tokens = len(tokenizer.encode(sentence))
                
                if sentence_tokens > max_tokens:
                    logger.warning("A single sentence is too long and will be split directly.")
                    words = sentence.split()
                    temp_sentence = ""
                    temp_tokens = 0
                    
                    for word in words:
                        word_tokens = len(tokenizer.encode(word))
                        if temp_tokens + word_tokens > max_tokens:
                            if temp_sentence:
                                chunks.append(temp_sentence.strip())
                            temp_sentence = word + " "
                            temp_tokens = word_tokens
                        else:
                            temp_sentence += word + " "
                            temp_tokens += word_tokens
                    
                    if temp_sentence:
                        chunks.append(temp_sentence.strip())
                    continue

                if current_tokens + sentence_tokens > max_tokens:
                    if current_chunk:
                        chunks.append(current_chunk.strip())
                    current_chunk = sentence + " "
                    current_tokens = sentence_tokens
                else:
                    current_chunk += sentence + " "
                    current_tokens += sentence_tokens

            if current_chunk:
                chunks.append(current_chunk.strip())

            return chunks
        except Exception as e:
            logger.error(f"Error splitting text into chunks: {str(e)}")
            return [text]  
    def get_models(self) -> Tuple[Any, Any]:
        """Initialize and return LLM models"""
        try:
            llm = ChatGroq(
                model_name="meta-llama/llama-4-scout-17b-16e-instruct",
                api_key=GROQ_API_KEY
            )
            
            genai.configure(api_key=google_Api_key)
            model = genai.GenerativeModel("gemini-1.5-flash")
            
            return llm, model
        except Exception as e:
            logger.error(f"Error initializing models: {str(e)}")
            raise Exception(f"Failed to initialize models: {str(e)}")

    def text_summary_chain(self, llm) -> Any:
        """Create text summarization chain"""
        prompt = PromptTemplate.from_template("""
You are an intelligent and concise summarization assistant.

A document was provided by the user and its content—including sections like the index, cover page, and references—was extracted and split into multiple text chunks.

Your task is to summarize **each individual chunk**. The summary should:
- Be factual, clear, and concise
- It should be of high quality no useful knowledge should be missed.
- Highlight important definitions, processes, relationships, or data
- Avoid irrelevant details like page numbers, boilerplate text, or formatting artifacts
- Do not start with "this is summary" or "summary as heading"....etc. I want only direct summary.

Chunk:
{chunk}

Return only the summary, focusing on making it informative and quiz-relevant.
""")
        return prompt | llm | StrOutputParser()

    def get_final_summaries(self, filepath: str, llm) -> str:
        """Generate final summary from document"""
        try:
            text_content, image_content = self.get_and_extract_file(filepath)
            
            if not text_content.strip():
                logger.warning("No text content found in document")
                return ""
            
            chunks = self.split_text_into_chunks(text_content, model_name="gpt-3.5-turbo", max_tokens=1000)
            
            if not chunks:
                return ""
            
            chain = self.text_summary_chain(llm)
            all_summaries = []
            
            for chunk in chunks:
                try:
                    summary = chain.invoke({"chunk": chunk})
                    if summary:
                        all_summaries.append(summary)
                except Exception as e:
                    logger.error(f"Error generating summary for chunk: {str(e)}")
                    continue
            
            final_summary = "\n\n".join(all_summaries)
            return final_summary
        except Exception as e:
            logger.error(f"Error generating final summaries: {str(e)}")
            return ""

    def describe_image_base64(self, image_b64: str) -> str:
        """Describe image using base64 encoding"""
        for attempt in range(self.max_retries):
            try:
                image_data = base64.b64decode(image_b64)
                image = Image.open(BytesIO(image_data))
                llm, model = self.get_models()
                
                prompt = (
                    "You are an AI assistant. Give a short 3-line educational description of this image. "
                    "Focus on what it visually represents, such as diagrams, graphs, labeled content, or learning concepts."
                )

                response = model.generate_content([prompt, image])

                if hasattr(response, 'text') and response.text and "Error:" not in response.text:
                    return response.text.strip()
                else:
                    logger.warning("Unexpected content or format issue in image description")
                    return ""

            except Exception as e:
                if "429" in str(e) or "quota" in str(e).lower():
                    logger.warning(f"Rate limit reached. Waiting {self.wait_time} seconds before retry...")
                    time.sleep(self.wait_time)
                else:
                    logger.error(f"Error describing image: {str(e)}")
                    return ""

        return ""

    def get_image_summary(self, filepath: str) -> str:
        """Generate summary of all images in document"""
        try:
            text_content, image_content = self.get_and_extract_file(filepath)
            
            if not image_content:
                logger.info("No images found in document")
                return ""
            
            image_summaries = []
            
            for i, img_b64 in enumerate(image_content):
                if i >= 50:  
                    break
                    
                summary = self.describe_image_base64(img_b64)
                if summary:
                    image_summaries.append(summary)
            
            final_image_summary = "\n\n".join(image_summaries)
            return final_image_summary
        except Exception as e:
            logger.error(f"Error generating image summary: {str(e)}")
            return ""

    def count_tokens(self, text: str, model: str = "gpt-3.5-turbo") -> int:
        """Count tokens in text"""
        try:
            tokenizer = tiktoken.encoding_for_model(model)
            tokens = tokenizer.encode(text)
            return len(tokens)
        except Exception as e:
            logger.error(f"Error counting tokens: {str(e)}")
            return len(text.split())  

    def get_text_quiz_prompt(self) -> str:
        """Get prompt template for text-based quiz generation"""
        return """You are an expert computer science educator and quiz designer.

Generate a set of exactly **{no_of_questions} multiple-choice questions** (MCQs) for students who uploads a book and the summary of the book is given to you.
Book Summary:
{summary}

Follow these strict instructions:

1. Cover a wide range of important **subtopics** from the given topic in interview level.
2. All the questions should be generated from given summary only. Hints can be generated by your own.
3. Do not ask simple and silly questions. Questions standard should be very high and practical.
4. The questions should be arranged by **difficulty level**:
  - 2 Medium questions (basic concepts and definitions)
  - 4 Medium to Advanced questions (application-level, comparative questions)
  - 4 Advanced questions (analysis, multi-step reasoning, real-world problem-solving)
5. For each question, provide:
  - The question text
  - Exactly 4 answer options (labeled A, B, C, D)
  - The correct option label
  - A brief explanation of the correct answer
6. Generate random questions every time.
7. Focus on asking programming and coding questions in case of programming related topics.
8. Format your output strictly in this JSON structure:
```json
[
{{
  "level": "medium",
  "question": "Your question here",
  "options": {{
    "A": "Option A",
    "B": "Option B",
    "C": "Option C",
    "D": "Option D"
  }},
  "answer": "C",
  "hint": "A basic hint here",
  "explanation": "Your explanation here"
}},
...
]
```"""

    def get_image_quiz_prompt(self) -> str:
        """Get prompt template for image-based quiz generation"""
        return """You are an expert computer science educator and quiz designer.

Generate a set of exactly **3 multiple-choice questions** (MCQs) for students who uploads a book and the summary of the images in the book is given to you.
The summary is generated in such a way that one by one image is given to the gemini flash 1.5 then it gives a 3 line summary for each image
Images Summary:
{img_summary}

Follow these strict instructions:

1. Cover a wide range of important **subtopics** from the given topic in interview level.
2. All the questions should be generated from given summary only. Hints can be generated by your own.
3. Do not ask simple and silly questions. Questions standard should be very high and practical.
4. The questions should be arranged by **difficulty level**:
  - 3 Advanced questions (analysis, multi-step reasoning, real-world problem-solving)
5. For each question, provide:
  - The question text
  - Exactly 4 answer options (labeled A, B, C, D)
  - The correct option label
  - A brief explanation of the correct answer
6. Generate random questions every time.
7. Focus on asking programming and coding questions in case of programming related topics.
8. Format your output strictly in this JSON structure:
```json
[
{{
  "level": "advanced",
  "question": "Your question here",
  "options": {{
    "A": "Option A",
    "B": "Option B",
    "C": "Option C",
    "D": "Option D"
  }},
  "answer": "C",
  "hint": "A basic hint here",
  "explanation": "Your explanation here"
}},
...
]
```"""

    def cleaner(self, llm_response: str) -> List[Dict[str, Any]]:
        """Clean and parse LLM response to JSON"""
        try:
            cleaned = re.sub(r"```json|```", "", llm_response).strip()
            answer = json.loads(cleaned)
            return answer
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing error: {str(e)}")
            logger.error(f"Response content: {llm_response}")
            raise Exception(f"Failed to parse quiz response: {str(e)}")

    def get_quiz_llm(self) -> Any:
        """Get LLM model for quiz generation"""
        try:
            model = ChatGroq(
                model_name="llama-3.3-70b-versatile",
                temperature=0.5,
                groq_api_key=GROQ_API_KEY
            )
            return model
        except Exception as e:
            logger.error(f"Error initializing quiz LLM: {str(e)}")
            raise Exception(f"Failed to initialize quiz LLM: {str(e)}")

    async def quiz_with_text(self, llm, final_summary: str) -> List[Dict[str, Any]]:
        """Generate quiz questions from text summary"""
        if not final_summary.strip():
            return []
        
        try:
            results = []
            sys_prompt = self.get_text_quiz_prompt()
            num_chunks = self.count_tokens(final_summary) // 20000
            
            if num_chunks <= 1:
                prompt = sys_prompt.format(no_of_questions=10, summary=final_summary)
                result = await llm.ainvoke(prompt)
                final_ans = self.cleaner(result.content)
                return final_ans
            else:
                questions_per_chunk = max(3, 10 // max(2, num_chunks))
                token_limit = self.count_tokens(final_summary) // max(2, num_chunks)
                
                chunks = self.split_text_into_chunks(final_summary, model_name="gpt-3.5-turbo", max_tokens=token_limit)
                
                for chunk in chunks:
                    prompt = sys_prompt.format(no_of_questions=questions_per_chunk, summary=chunk)
                    result = await llm.ainvoke(prompt)
                    chunk_questions = self.cleaner(result.content)
                    results.extend(chunk_questions)
                
                return results[:10]  # Limit to 10 questions
        except Exception as e:
            logger.error(f"Error generating text quiz: {str(e)}")
            return []

    async def quiz_with_images(self, img_summary: str, llm) -> List[Dict[str, Any]]:
        """Generate quiz questions from image summary"""
        if not img_summary.strip():
            return []
        
        try:
            sys_prompt = self.get_image_quiz_prompt()
            prompt = sys_prompt.format(img_summary=img_summary)
            result = await llm.ainvoke(prompt)
            final_ans = self.cleaner(result.content)
            return final_ans
        except Exception as e:
            logger.error(f"Error generating image quiz: {str(e)}")
            return []

    async def quiz_with_file(self, filepath: str) -> Dict[str, Any]:
        """Main function to generate quiz from file (PDF, DOCX, or PPTX)"""
        try:
            # Initialize models
            text_llm, _ = self.get_models()
            quiz_llm = self.get_quiz_llm()
            
            # Generate text summary
            logger.info(f"Generating text summary from {filepath}")
            final_summary = self.get_final_summaries(filepath, text_llm)
            
            # Generate image summary
            logger.info(f"Generating image summary from {filepath}")
            img_summary = self.get_image_summary(filepath)
            logger.info("Generating quiz questions")
            text_quiz_task = self.quiz_with_text(quiz_llm, final_summary)
            image_quiz_task = self.quiz_with_images(img_summary, quiz_llm)
            
            text_quiz, image_quiz = await asyncio.gather(text_quiz_task, image_quiz_task)
            
            # Combine results
            combined_quiz = []
            if text_quiz:
                combined_quiz.extend(text_quiz)
            if image_quiz:
                combined_quiz.extend(image_quiz)
            
            if not combined_quiz:
                return {
                    "success": False,
                    "message": "No quiz questions could be generated from the provided file",
                    "data": []
                }
            
            return {
                "success": True,
                "message": f"Successfully generated {len(combined_quiz)} quiz questions",
                "data": combined_quiz,
                "metadata": {
                    "text_questions": len(text_quiz) if text_quiz else 0,
                    "image_questions": len(image_quiz) if image_quiz else 0,
                    "total_questions": len(combined_quiz),
                    "file_type": Path(filepath).suffix.lower()
                }
            }
        
        except Exception as e:
            logger.error(f"Error in quiz_with_file: {str(e)}")
            return {
                "success": False,
                "message": f"Error generating quiz: {str(e)}",
                "data": []
            }
quiz_service = QuizGenerationService()